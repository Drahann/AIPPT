#!/usr/bin/env python3
"""extract_deck — build the deck_record (PPT-level skin) for one template.

Sources, per ARCH §2.0:
  - ppt/theme/theme1.xml  -> clrScheme (12 colors) + fontScheme (major/minor)
  - measured slides       -> the REAL per-slot fonts/sizes (the theme font is often
                             just 微软雅黑 while slides override with display fonts),
                             aggregated into a size_ramp + family roles.
Output: <out>/deck_record.json

Usage: python extract_deck.py <pptx> <out_dir>
"""
import os, sys, json, zipfile, collections
import xml.etree.ElementTree as ET
from pptx import Presentation

A = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
src, outdir = sys.argv[1], sys.argv[2]
os.makedirs(outdir, exist_ok=True)

# ---- theme palette + font scheme ----
z = zipfile.ZipFile(src)
root = ET.fromstring(z.read('ppt/theme/theme1.xml').decode('utf-8'))
cs = root.find('.//' + A + 'clrScheme')
palette = {}
for c in cs:
    tag = c.tag.split('}')[-1]
    srgb, sysc = c.find(A + 'srgbClr'), c.find(A + 'sysClr')
    palette[tag] = '#' + (srgb.get('val') if srgb is not None else sysc.get('lastClr'))
fs = root.find('.//' + A + 'fontScheme')
def face(grp):
    g = fs.find(A + grp)
    return g.find(A + 'latin').get('typeface'), (g.find(A + 'ea').get('typeface') if g.find(A + 'ea') is not None else '')
maj_latin, maj_ea = face('majorFont')
min_latin, min_ea = face('minorFont')

# ---- measured fonts/sizes across slides (the real display fonts) ----
prs = Presentation(src)
fam_by_size = collections.Counter()    # (font, px_bucket) -> count
fam_count = collections.Counter()
sizes = collections.Counter()
for sl in prs.slides:
    for sh in sl.shapes:
        if not getattr(sh, 'has_text_frame', False):
            continue
        for para in sh.text_frame.paragraphs:
            for run in para.runs:
                f = run.font
                rpr = run._r.find(A + 'rPr')
                nm = f.name
                if nm is None and rpr is not None:
                    for tag in ('latin', 'ea'):
                        e = rpr.find(A + tag)
                        if e is not None and e.get('typeface'):
                            nm = e.get('typeface'); break
                if nm and not nm.startswith('+'):
                    fam_count[nm] += 1
                    if f.size is not None:
                        px = round(f.size.pt * 4 / 3)
                        fam_by_size[(nm, px)] += 1
                        sizes[px] += 1

top_fams = [f for f, _ in fam_count.most_common(8)]
# title family = the most common LARGE-size display font; body = most common small font
def family_at(min_px, max_px):
    cand = collections.Counter()
    for (nm, px), n in fam_by_size.items():
        if min_px <= px <= max_px:
            cand[nm] += n
    return cand.most_common(1)[0][0] if cand else (top_fams[0] if top_fams else '')

deck = {
    'deck_id': os.path.splitext(os.path.basename(src))[0],
    'canvas': [1280, 720],
    'theme': {
        'palette': palette,                 # raw 12-color clrScheme
        'font_scheme': {'major': [maj_latin, maj_ea], 'minor': [min_latin, min_ea]},
    },
    'typography': {                          # MEASURED display fonts (authoritative for slots)
        'title_family': family_at(40, 200),
        'subtitle_family': family_at(28, 40),
        'body_family': family_at(10, 24),
        'families_seen': top_fams,
        'size_ramp': [s for s, _ in sizes.most_common(10)],
    },
    # semantic palette aliases used by resolve(page, skin); measured/deck-known values
    'skin': {
        'primary': palette.get('accent2', '#3F81F6'),
        'accent_cyan': palette.get('accent3', '#5AF3FD'),
        'bg_deep': palette.get('dk1', '#030E22'),
        'text': palette.get('lt1', '#FFFFFF'),
        'text_dim': palette.get('lt2', '#E7E6E6'),
    },
}
json.dump(deck, open(os.path.join(outdir, 'deck_record.json'), 'w', encoding='utf-8'),
          ensure_ascii=False, indent=1)
print(json.dumps(deck, ensure_ascii=False, indent=1))
