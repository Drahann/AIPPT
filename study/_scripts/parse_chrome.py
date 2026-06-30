#!/usr/bin/env python3
"""Dump a slide's LAYOUT + MASTER persistent chrome (images, boxes, background).
The top ornament / full-bleed bg often live here, not in slide shapes.

Usage: python parse_chrome.py <pptx> <slide_no_1based> <out_assets_dir>
"""
import os, sys
from pptx import Presentation

A = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
P = '{http://schemas.openxmlformats.org/presentationml/2006/main}'
R = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}'

src, sno, outdir = sys.argv[1], int(sys.argv[2]), sys.argv[3]
os.makedirs(outdir, exist_ok=True)
prs = Presentation(src)
W, H = int(prs.slide_width), int(prs.slide_height)
SX, SY = 1280.0 / W, 720.0 / H
slide = prs.slides[sno - 1]
layout = slide.slide_layout
master = layout.slide_master
lines, imgn = [], [0]


def grp_xfrm(el):
    xf = el.find(P + 'grpSpPr/' + A + 'xfrm')
    o, e = xf.find(A + 'off'), xf.find(A + 'ext')
    co, ce = xf.find(A + 'chOff'), xf.find(A + 'chExt')
    return (int(o.get('x')), int(o.get('y')), int(e.get('cx')), int(e.get('cy')),
            int(co.get('x')), int(co.get('y')), int(ce.get('cx')), int(ce.get('cy')))


def walk(shapes, tag, part, mapxy, kx, ky):
    for sh in shapes:
        try:
            l, t, w, h = sh.left, sh.top, sh.width, sh.height
        except Exception:
            l = None
        if sh.shape_type == 6:
            ox, oy, ex, ey, cox, coy, cex, cey = grp_xfrm(sh._element)
            gx0, gy0 = mapxy(ox, oy)
            walk(sh.shapes, tag, part,
                 (lambda x, y, gx0=gx0, gy0=gy0, cox=cox, coy=coy, nx=kx*(ex/cex if cex else 1), ny=ky*(ey/cey if cey else 1): (gx0+(x-cox)*nx, gy0+(y-coy)*ny)),
                 kx*(ex/cex if cex else 1), ky*(ey/cey if cey else 1))
            continue
        if l is None:
            continue
        X, Y = mapxy(l, t)
        box = [round(X*SX, 1), round(Y*SY, 1), round(w*kx*SX, 1), round(h*ky*SY, 1)]
        extra = ''
        if sh.shape_type == 13:
            try:
                img = sh.image
                imgn[0] += 1
                fn = f"{part}{imgn[0]:02d}.{img.ext}"
                open(os.path.join(outdir, fn), 'wb').write(img.blob)
                extra = f" img={fn}({img.size[0]}x{img.size[1]})"
            except Exception as e:
                extra = f" img_err={str(e)[:30]}"
        try:
            if sh.has_text_frame and sh.text_frame.text.strip():
                extra += f"  <<{sh.text_frame.text.strip()[:40]}>>"
        except Exception:
            pass
        lines.append(f"  {tag} [x{box[0]:>6} y{box[1]:>5} {box[2]:>6}x{box[3]:>5}] {sh.name[:18]:<18}{extra}")


def bg_of(part_el, label):
    bg = part_el.find(P + 'cSld/' + P + 'bg')
    if bg is None:
        return
    blip = bg.find('.//' + A + 'blip')
    if blip is not None:
        lines.append(f"  [BG-{label}] blipFill image (embedded bg) embed={blip.get(R+'embed')}")
    else:
        fill = list(bg.find(P + 'bgPr') if bg.find(P + 'bgPr') is not None else bg)
        lines.append(f"  [BG-{label}] {bg[0].tag.split('}')[-1] if len(bg) else '?'} (non-image fill)")


for el, lab, part in [(slide._element, 'SLIDE', 'S'), (layout._element, 'LAYOUT', 'L'), (master._element, 'MASTER', 'M')]:
    lines.append(f"--- {lab} ---")
    bg_of(el, lab)

imgn[0] = 0
lines.append("--- LAYOUT shapes ---")
walk(layout.shapes, 'L', 'L', (lambda x, y: (x, y)), 1.0, 1.0)
imgn[0] = 0
lines.append("--- MASTER shapes ---")
walk(master.shapes, 'M', 'M', (lambda x, y: (x, y)), 1.0, 1.0)

rep = "\n".join(lines)
open(os.path.join(outdir, '_chrome.txt'), 'w', encoding='utf-8').write(rep)
try:
    print(rep)
except Exception:
    print("chrome report ->", outdir + "/_chrome.txt")
