# -*- coding: utf-8 -*-
"""Structural/quantitative extraction over a sample CSV (python-pptx + Pillow).
Measures: aspect, density (shapes/pics per slide), image-area ratio, full-bleed bg rate,
REAL palette (hard-coded srgb fills + dominant colors of full-bleed bg images), live-text fonts,
cover composition. Writes per-deck + per-family JSON and a markdown summary.
Usage: python metrics_extract.py <sample_csv> <out_prefix>"""
import os, sys, io, csv, json, collections
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_FILL
from PIL import Image

sample_csv = sys.argv[1] if len(sys.argv)>1 else r"W:\ppt\study\_scripts\sample_R1.csv"
out_prefix = sys.argv[2] if len(sys.argv)>2 else r"W:\ppt\study\metrics\R1"

def iter_shapes(shapes):
    for sh in shapes:
        yield sh
        try:
            if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from iter_shapes(sh.shapes)
        except: pass

def fill_rgb(sh):
    try:
        f = sh.fill
        if f.type == MSO_FILL.SOLID:
            return str(f.fore_color.rgb)   # 'RRGGBB'; raises if theme/scheme color
    except: pass
    return None

def line_rgb(sh):
    try:
        ln = sh.line
        return str(ln.color.rgb)
    except: pass
    return None

def pic_dominant(sh):
    try:
        blob = sh.image.blob
        im = Image.open(io.BytesIO(blob)).convert("RGB").resize((40,40))
        q = im.quantize(colors=5)
        pal = q.getpalette()
        cnt = collections.Counter(q.getdata())
        idx = cnt.most_common(1)[0][0]
        return '%02X%02X%02X' % (pal[idx*3], pal[idx*3+1], pal[idx*3+2])
    except: return None

rows = list(csv.DictReader(open(sample_csv, encoding="utf-8-sig")))
per_deck = []
fam_agg = collections.defaultdict(lambda: {
    "decks":0, "aspect":collections.Counter(), "shapes_per_slide":[], "pics_per_slide":[],
    "img_area_ratio":[], "fullbleed_rate":[], "fills":collections.Counter(),
    "lines":collections.Counter(), "bg_dom":collections.Counter(), "fonts":collections.Counter(),
    "cover_img_ratio":[], "cover_live_title":0, "n_slides":[],
})

for r in rows:
    did=r["deck_id"]; fam=r["family"]; path=r["fullpath"]
    try:
        prs = Presentation(path)
    except Exception as e:
        print(f"  ERR open {did}: {e}"); continue
    SW, SH = prs.slide_width, prs.slide_height
    area = SW*SH if SW and SH else 1
    aspect = round(SW/SH,3) if SH else 0
    asp_label = "16:9" if abs(aspect-16/9)<0.03 else ("4:3" if abs(aspect-4/3)<0.03 else f"{aspect}")
    d = {"deck_id":did,"family":fam,"n_slides":len(prs.slides),"aspect":asp_label,
         "shapes_per_slide":0,"pics_per_slide":0,"img_area_ratio":0,"fullbleed_rate":0}
    A=fam_agg[fam]; A["decks"]+=1; A["aspect"][asp_label]+=1; A["n_slides"].append(len(prs.slides))
    tot_shapes=tot_pics=0; img_ratios=[]; fullbleeds=0
    for si,slide in enumerate(prs.slides):
        shapes=list(iter_shapes(slide.shapes))
        ns=len(shapes); npic=0; pic_area=0; has_fullbleed=False
        for sh in shapes:
            try: stype=sh.shape_type
            except: stype=None
            # fills/lines colors
            fr=fill_rgb(sh)
            if fr and fr!="000000" and fr!="FFFFFF": A["fills"][fr]+=1
            lr=line_rgb(sh)
            if lr and lr not in ("000000","FFFFFF"): A["lines"][lr]+=1
            # fonts
            if sh.has_text_frame:
                for para in sh.text_frame.paragraphs:
                    for run in para.runs:
                        fn=run.font.name
                        if fn: A["fonts"][fn]+=1
            if stype==MSO_SHAPE_TYPE.PICTURE:
                npic+=1
                try:
                    a=sh.width*sh.height; pic_area+=a
                    if a>=0.82*area and sh.width>=0.9*SW: has_fullbleed=True
                except: pass
                if si<2 or has_fullbleed:
                    dom=pic_dominant(sh)
                    if dom: A["bg_dom"][dom]+=1
        tot_shapes+=ns; tot_pics+=npic
        img_ratios.append(min(1.0, pic_area/area))
        if has_fullbleed: fullbleeds+=1
        if si==0:
            cover_ratio=min(1.0,pic_area/area); A["cover_img_ratio"].append(cover_ratio)
            live_title=any(sh.has_text_frame and len(sh.text_frame.text.strip())>=4
                           and any(run.font.size and run.font.size.pt>=28
                                   for p in sh.text_frame.paragraphs for run in p.runs)
                           for sh in shapes)
            if live_title: A["cover_live_title"]+=1
    nsl=max(1,len(prs.slides))
    d["shapes_per_slide"]=round(tot_shapes/nsl,1)
    d["pics_per_slide"]=round(tot_pics/nsl,1)
    d["img_area_ratio"]=round(sum(img_ratios)/nsl,3)
    d["fullbleed_rate"]=round(fullbleeds/nsl,3)
    A["shapes_per_slide"].append(d["shapes_per_slide"]); A["pics_per_slide"].append(d["pics_per_slide"])
    A["img_area_ratio"].append(d["img_area_ratio"]); A["fullbleed_rate"].append(d["fullbleed_rate"])
    per_deck.append(d)
    print(f"  {did:24s} {asp_label:5s} slides={d['n_slides']:2d} sh/s={d['shapes_per_slide']:4} pic/s={d['pics_per_slide']:4} imgR={d['img_area_ratio']:.2f} fb={d['fullbleed_rate']:.2f}")

def avg(x): return round(sum(x)/len(x),2) if x else 0
fam_out={}
for fam,A in fam_agg.items():
    fam_out[fam]={
        "decks":A["decks"], "aspect":dict(A["aspect"]),
        "avg_slides":avg(A["n_slides"]),
        "avg_shapes_per_slide":avg(A["shapes_per_slide"]),
        "avg_pics_per_slide":avg(A["pics_per_slide"]),
        "avg_img_area_ratio":avg(A["img_area_ratio"]),
        "avg_fullbleed_rate":avg(A["fullbleed_rate"]),
        "cover_img_ratio":avg(A["cover_img_ratio"]),
        "cover_live_title":f"{A['cover_live_title']}/{A['decks']}",
        "top_fills":A["fills"].most_common(12),
        "top_lines":A["lines"].most_common(8),
        "top_bg_dominant":A["bg_dom"].most_common(12),
        "top_fonts":A["fonts"].most_common(12),
    }

json.dump({"per_deck":per_deck,"per_family":fam_out},
          open(out_prefix+"_metrics.json","w",encoding="utf-8"), ensure_ascii=False, indent=1)

# markdown summary
L=["# Metrics — "+os.path.basename(sample_csv),""]
for fam,o in fam_out.items():
    L.append(f"\n## {fam}  ({o['decks']} decks)")
    L.append(f"- aspect: {o['aspect']} | avg slides {o['avg_slides']}")
    L.append(f"- density: shapes/slide **{o['avg_shapes_per_slide']}**, pics/slide **{o['avg_pics_per_slide']}**, img-area **{o['avg_img_area_ratio']}**, full-bleed rate **{o['avg_fullbleed_rate']}**")
    L.append(f"- cover: img-area **{o['cover_img_ratio']}**, live-title decks {o['cover_live_title']}")
    L.append(f"- top hard-coded fills: {', '.join('#'+h+'×'+str(n) for h,n in o['top_fills'])}")
    L.append(f"- top bg dominant: {', '.join('#'+h+'×'+str(n) for h,n in o['top_bg_dominant'])}")
    L.append(f"- top fonts: {', '.join(f+'×'+str(n) for f,n in o['top_fonts'])}")
open(out_prefix+"_metrics.md","w",encoding="utf-8").write("\n".join(L))
print("\n-> "+out_prefix+"_metrics.json / .md")