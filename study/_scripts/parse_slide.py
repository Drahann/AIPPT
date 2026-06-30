#!/usr/bin/env python3
"""Parse ONE pptx slide's real layout: every shape's absolute px box (with group
transform), text, and extract every embedded picture. So replication uses REAL
image positions/sizes, not eyeballed guesses.

Usage: python parse_slide.py <pptx> <slide_no_1based> <out_assets_dir>
"""
import os, sys, json
from pptx import Presentation

A = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
P = '{http://schemas.openxmlformats.org/presentationml/2006/main}'

src, sno, outdir = sys.argv[1], int(sys.argv[2]), sys.argv[3]
os.makedirs(outdir, exist_ok=True)
prs = Presentation(src)
W, H = int(prs.slide_width), int(prs.slide_height)
SX, SY = 1280.0 / W, 720.0 / H
slide = prs.slides[sno - 1]
rows = []
imgn = [0]


def grp_xfrm(el):
    xfrm = el.find(P + 'grpSpPr/' + A + 'xfrm')
    off, ext = xfrm.find(A + 'off'), xfrm.find(A + 'ext')
    cho, che = xfrm.find(A + 'chOff'), xfrm.find(A + 'chExt')
    return (int(off.get('x')), int(off.get('y')), int(ext.get('cx')), int(ext.get('cy')),
            int(cho.get('x')), int(cho.get('y')), int(che.get('cx')), int(che.get('cy')))


def walk(shapes, mapxy, kx, ky):
    for sh in shapes:
        try:
            l, t, w, h = sh.left, sh.top, sh.width, sh.height
        except Exception:
            l = t = w = h = None
        if sh.shape_type == 6:  # GROUP
            ox, oy, ex, ey, cox, coy, cex, cey = grp_xfrm(sh._element)
            gx0, gy0 = mapxy(ox, oy)
            nkx, nky = kx * (ex / cex if cex else 1), ky * (ey / cey if cey else 1)
            walk(sh.shapes,
                 (lambda x, y, gx0=gx0, gy0=gy0, cox=cox, coy=coy, nkx=nkx, nky=nky:
                  (gx0 + (x - cox) * nkx, gy0 + (y - coy) * nky)),
                 nkx, nky)
            continue
        if l is None:
            continue
        X, Y = mapxy(l, t)
        box = [round(X * SX, 1), round(Y * SY, 1), round(w * kx * SX, 1), round(h * ky * SY, 1)]
        rec = {'name': sh.name, 'type': str(sh.shape_type), 'box': box}
        # rotation + flip (from xfrm)
        xf = sh._element.find('.//' + A + 'xfrm')
        if xf is not None:
            rot = int(xf.get('rot', '0'))
            if rot:
                rec['rot'] = round(rot / 60000.0, 1)
            if xf.get('flipH') in ('1', 'true'):
                rec['flipH'] = True
            if xf.get('flipV') in ('1', 'true'):
                rec['flipV'] = True
        try:
            if sh.has_text_frame and sh.text_frame.text.strip():
                tf = sh.text_frame
                rec['text'] = tf.text.strip().replace('\n', ' / ')[:70]
                # first run font props
                for para in tf.paragraphs:
                    if para.alignment is not None:
                        rec['align'] = str(para.alignment).split()[0]
                    for run in para.runs:
                        f = run.font
                        sz = f.size.pt if f.size is not None else None
                        nm2 = f.name
                        # fall back to rPr latin/ea typeface
                        rpr = run._r.find(A + 'rPr')
                        if nm2 is None and rpr is not None:
                            for tag in ('latin', 'ea'):
                                el = rpr.find(A + tag)
                                if el is not None and el.get('typeface'):
                                    nm2 = el.get('typeface'); break
                        col = None
                        if rpr is not None:
                            sf = rpr.find(A + 'solidFill/' + A + 'srgbClr')
                            if sf is not None:
                                col = '#' + sf.get('val')
                        if sz:
                            rec['font_pt'] = round(sz, 1)
                            rec['font_px'] = round(sz * 4 / 3, 1)
                        if nm2:
                            rec['font'] = nm2
                        if f.bold:
                            rec['bold'] = True
                        if col:
                            rec['color'] = col
                        break
                    if 'font_pt' in rec or 'font' in rec:
                        break
        except Exception:
            pass
        if sh.shape_type == 13:  # PICTURE
            try:
                img = sh.image
                imgn[0] += 1
                fn = f"pic{imgn[0]:02d}.{img.ext}"
                open(os.path.join(outdir, fn), 'wb').write(img.blob)
                rec['image'] = fn
                rec['img_native'] = f"{img.size[0]}x{img.size[1]}"
            except Exception as e:
                rec['image_err'] = str(e)[:40]
        rows.append(rec)


walk(slide.shapes, (lambda x, y: (x, y)), 1.0, 1.0)
lines = [f"slide {sno} of {len(prs.slides)} | canvas EMU {W}x{H} -> 1280x720 | leaf shapes: {len(rows)} | images: {imgn[0]}", "=" * 100]
for i, r in enumerate(rows):
    bx = r['box']
    tag = 'PIC ' if 'image' in r else ('TXT ' if 'text' in r else 'shp ')
    fl = ''
    if r.get('rot'):
        fl += f" rot{r['rot']}"
    if r.get('flipH'):
        fl += ' flipH'
    if r.get('flipV'):
        fl += ' flipV'
    extra = f" img={r['image']}({r['img_native']})" if 'image' in r else ''
    if 'text' in r:
        fnt = []
        if 'font' in r:
            fnt.append(r['font'])
        if 'font_px' in r:
            fnt.append(f"{r['font_px']}px/{r['font_pt']}pt")
        if r.get('bold'):
            fnt.append('B')
        if 'color' in r:
            fnt.append(r['color'])
        if 'align' in r:
            fnt.append(r['align'])
        extra += f"  {{{' '.join(fnt)}}}  <<{r['text']}>>"
    lines.append(f"{i:>2} {tag}[x{bx[0]:>6} y{bx[1]:>5}  {bx[2]:>6}x{bx[3]:>5}]{fl}  {r['name'][:18]:<18}{extra}")
report = "\n".join(lines)
open(os.path.join(outdir, '_report.txt'), 'w', encoding='utf-8').write(report)
json.dump({'canvas': [1280, 720], 'shapes': rows}, open(os.path.join(outdir, '_layout.json'), 'w', encoding='utf-8'), ensure_ascii=False, indent=1)
try:
    print(report)
except Exception:
    print(f"report written ({len(rows)} shapes, {imgn[0]} images) -> {outdir}/_report.txt")
