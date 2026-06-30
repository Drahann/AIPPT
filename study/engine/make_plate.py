#!/usr/bin/env python3
"""make_plate — render a faithful BASE PLATE: the slide with the shapes I'm going to
redraw removed, everything else (chrome + decorations + kept images) baked by COM.

Strategy (architecture: reuse-verbatim chrome/decoration, redraw only what changes):
copy pptx -> on the target slide, delete the shapes flagged for redraw (text / a hero
picture by native size / native charts / named shapes) -> COM-render that slide. What
remains is the exact visual base PowerPoint draws; the renderer overlays new text,
content-regen images, and vector charts on top.

Usage:
  python make_plate.py <pptx> <slide_1based> <out_png>
      [--kill-text] [--kill-pic-native 1024x1024,800x800]
      [--kill-charts] [--kill-names "名字1,名字2"] [--keep-names "n1,n2"]
"""
import os, sys, shutil, argparse, tempfile
from pptx import Presentation
import win32com.client

A = '{http://schemas.openxmlformats.org/drawingml/2006/main}'

ap = argparse.ArgumentParser()
ap.add_argument('pptx'); ap.add_argument('slide', type=int); ap.add_argument('out')
ap.add_argument('--kill-text', action='store_true')
ap.add_argument('--kill-pic-native', default='')
ap.add_argument('--kill-charts', action='store_true')
ap.add_argument('--kill-names', default='')
ap.add_argument('--keep-names', default='')
ap.add_argument('--kill-boxsize', default='', help='encoding-safe: rendered box WxH, e.g. 188x162,33x74')
ap.add_argument('--kill-region', default='', help='x0,y0,x1,y1 ; kill shapes whose box center is inside')
a = ap.parse_args()

kill_native = set(x for x in a.kill_pic_native.split(',') if x)
kill_names = set(x for x in a.kill_names.split(',') if x)
keep_names = set(x for x in a.keep_names.split(',') if x)
kill_boxsize = set(x for x in a.kill_boxsize.split(',') if x)
kill_region = [float(v) for v in a.kill_region.split(',')] if a.kill_region else None

W_, H_ = int(Presentation(a.pptx).slide_width), int(Presentation(a.pptx).slide_height)
SX_, SY_ = 1280.0 / W_, 720.0 / H_


GA = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
GP = '{http://schemas.openxmlformats.org/presentationml/2006/main}'


def grp_xfrm(el):
    xf = el.find(GP + 'grpSpPr/' + GA + 'xfrm')
    o, e = xf.find(GA + 'off'), xf.find(GA + 'ext')
    co, ce = xf.find(GA + 'chOff'), xf.find(GA + 'chExt')
    return (int(o.get('x')), int(o.get('y')), int(e.get('cx')), int(e.get('cy')),
            int(co.get('x')), int(co.get('y')), int(ce.get('cx')), int(ce.get('cy')))


def abs_box(sh, mapxy, kx, ky):
    """absolute 1280x720 box [x,y,w,h] honoring the group-transform chain."""
    try:
        X, Y = mapxy(sh.left, sh.top)
        return [X*SX_, Y*SY_, sh.width*kx*SX_, sh.height*ky*SY_]
    except Exception:
        return None

tmp = os.path.join(tempfile.gettempdir(), f"_plate_{os.getpid()}.pptx")
shutil.copy(a.pptx, tmp)
prs = Presentation(tmp)
slide = prs.slides[a.slide - 1]

removed = [0]


def consider(sh, box):
    """return True if this shape should be removed (redrawn later)."""
    if sh.name in keep_names:
        return False
    if sh.name in kill_names:
        return True
    if a.kill_charts and sh.has_chart:
        return True
    if a.kill_text:
        try:
            if sh.has_text_frame and sh.text_frame.text.strip():
                return True
        except Exception:
            pass
    if kill_native and sh.shape_type == 13:
        try:
            img = sh.image
            if f"{img.size[0]}x{img.size[1]}" in kill_native:
                return True
        except Exception:
            pass
    if box is not None:
        if kill_boxsize and f"{round(box[2])}x{round(box[3])}" in kill_boxsize:
            return True
        if kill_region:
            cx, cy = box[0] + box[2]/2, box[1] + box[3]/2
            if kill_region[0] <= cx <= kill_region[2] and kill_region[1] <= cy <= kill_region[3]:
                return True
    return False


def walk(shape_container, mapxy, kx, ky):
    for sh in list(shape_container.shapes):
        if sh.shape_type == 6:  # group: recurse with composed transform
            ox, oy, ex, ey, cox, coy, cex, cey = grp_xfrm(sh._element)
            gx0, gy0 = mapxy(ox, oy)
            nkx, nky = kx*(ex/cex if cex else 1), ky*(ey/cey if cey else 1)
            walk(sh, (lambda x, y, gx0=gx0, gy0=gy0, cox=cox, coy=coy, nkx=nkx, nky=nky:
                      (gx0+(x-cox)*nkx, gy0+(y-coy)*nky)), nkx, nky)
            continue
        if consider(sh, abs_box(sh, mapxy, kx, ky)):
            sh._element.getparent().remove(sh._element)
            removed[0] += 1


walk(slide, (lambda x, y: (x, y)), 1.0, 1.0)
prs.save(tmp)
print(f"removed {removed[0]} shapes; rendering plate...")

app = win32com.client.Dispatch("PowerPoint.Application")
pres = app.Presentations.Open(os.path.abspath(tmp), ReadOnly=True, WithWindow=False)
try:
    os.makedirs(os.path.dirname(os.path.abspath(a.out)), exist_ok=True)
    pres.Slides.Item(a.slide).Export(os.path.abspath(a.out), "PNG", 1280, 720)
    print("plate ->", a.out)
finally:
    pres.Close(); app.Quit()
    try:
        os.remove(tmp)
    except Exception:
        pass
