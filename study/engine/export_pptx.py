#!/usr/bin/env python3
"""export_pptx — assemble finished pages into a real 16:9 .pptx (S6 导出).

Each page in this pipeline renders to a deterministic, final 1280x720 image (page.png) from its
page.svg. The faithful, reliable export is therefore ONE full-bleed image per slide — pixel-identical
to what the render/gate approved (the SVG->native-shape path is parked, see PLAN_创赛特化改造.md). With
--hires the page's SVG is re-rasterized at 2x so the slide stays crisp on a projector.

Usage:
  python export_pptx.py -o deck.pptx <page_dir> [<page_dir> ...]
  python export_pptx.py -o deck.pptx --deck ../runs/rareearth        # auto-discover p*/ in page order
  python export_pptx.py -o deck.pptx --deck ../runs/rareearth --hires

A page_dir must contain page.png (and page.svg for --hires). Pages are ordered by the number in the
folder name (p18_财务规划 -> 18); pass dirs explicitly to force a custom order.
"""
import os, sys, re, glob, argparse, subprocess, tempfile

HERE = os.path.dirname(os.path.abspath(__file__))
EMU_PER_PX = 9525                         # 1280x720 px @96dpi == 13.333x7.5in; 1 px = 9525 EMU exactly
PT_PER_PX = 0.75                          # 1 px @96dpi == 0.75 pt

# alias / real font name -> a single installed family name for PowerPoint runs (heavy display vs body)
PPTX_FONT = {
    'DISP_TITLE': 'Alibaba PuHuiTi H', 'DISP_NUM': 'Alibaba PuHuiTi H', 'DISP_KAI': '华文行楷',
    'BODY': 'Microsoft YaHei', 'BODY_M': 'Microsoft YaHei', 'BODY_B': 'Microsoft YaHei',
    '阿里妈妈数黑体': 'Alibaba PuHuiTi H', '优设标题黑': 'Alibaba PuHuiTi H',
    '字体圈欣意冠黑体': '字体圈欣意冠黑体', '演示流云楷': '华文行楷',
    'OPPOSans R': 'Microsoft YaHei', 'OPPOSans M': 'Microsoft YaHei', 'OPPOSans B': 'Microsoft YaHei',
}
HEAVY = {'DISP_TITLE', 'DISP_NUM', 'BODY_B', '阿里妈妈数黑体', '优设标题黑', '字体圈欣意冠黑体'}


def pptx_font(name):
    return PPTX_FONT.get(name, name or 'Microsoft YaHei')


def mid_color(sh):
    """resolve a text shape's color to one hex (gradient -> middle stop) for an editable run."""
    if sh.get('grad'):
        stops = sh['grad']; return stops[len(stops) // 2].lstrip('#')
    return (sh.get('color') or '#FFFFFF').lstrip('#')


def page_num(d):
    m = re.search(r'p(\d+)', os.path.basename(os.path.normpath(d)))
    return int(m.group(1)) if m else 10**9


def collect(args):
    dirs = list(args.page_dirs)
    if args.deck:
        found = [d for d in glob.glob(os.path.join(args.deck, 'p*')) if os.path.isdir(d)]
        dirs += sorted(found, key=page_num)
    # de-dup preserving order
    seen, out = set(), []
    for d in dirs:
        ad = os.path.abspath(d)
        if ad not in seen and os.path.isdir(ad):
            seen.add(ad); out.append(ad)
    return out


def page_image(d, hires, tmp):
    """Return a path to this page's slide image (re-rasterizing the SVG at 2x if --hires)."""
    png = os.path.join(d, 'page.png')
    svg = os.path.join(d, 'page.svg')
    if hires and os.path.exists(svg):
        out = os.path.join(tmp, os.path.basename(d) + '@2x.png')
        r = subprocess.run([sys.executable, os.path.join(HERE, 'svg_to_png.py'), svg, out, '2560', '1440'],
                           capture_output=True, text=True)
        if os.path.exists(out):
            return out
        print(f"  ! hires re-raster failed for {os.path.basename(d)} ({r.stderr[-120:].strip()}); using page.png")
    return png if os.path.exists(png) else None


def textfree_bg(page_dir, tmp):
    """Re-rasterize page.svg with the content text (class="ct") removed, keeping plate + decorations
    + chart-internal labels, at 2x. This is the visual layer under the editable text boxes."""
    svg = os.path.join(page_dir, 'page.svg')
    if not os.path.exists(svg):
        return None
    lines = [ln for ln in open(svg, encoding='utf-8').read().split('\n') if 'class="ct"' not in ln]
    tsvg = os.path.join(tmp, os.path.basename(page_dir) + '_notext.svg')
    open(tsvg, 'w', encoding='utf-8').write('\n'.join(lines))
    out = os.path.join(tmp, os.path.basename(page_dir) + '_bg.png')
    subprocess.run([sys.executable, os.path.join(HERE, 'svg_to_png.py'), tsvg, out, '2560', '1440'],
                   capture_output=True, text=True)
    return out if os.path.exists(out) else None


def _style_run(run, sh, name, bold, hexc):
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    from pptx.oxml.ns import qn
    run.text = run.text
    f = run.font
    f.size = Pt(round(sh.get('font_px', 18) * PT_PER_PX, 1))
    f.bold = bold
    f.name = name                                          # latin face
    rPr = run._r.get_or_add_rPr()                          # also set East-Asian face so CJK uses it
    for tag in ('a:ea', 'a:cs'):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {}); rPr.append(el)
        el.set('typeface', name)
    try:
        f.color.rgb = RGBColor.from_string(hexc)
    except Exception:
        pass


def add_text_boxes(slide, page_dir):
    """Add NATIVE, editable text boxes from page.record.json — same box/font/size/color/align as render."""
    import json
    from pptx.util import Emu
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    rec = json.load(open(os.path.join(page_dir, 'page.record.json'), encoding='utf-8'))
    align_map = {'CENTER': PP_ALIGN.CENTER, 'RIGHT': PP_ALIGN.RIGHT, 'LEFT': PP_ALIGN.LEFT}
    n = 0
    for sh in rec.get('shapes', []):
        if sh.get('kind') != 'text' or not str(sh.get('text', '')).strip():
            continue
        x, y, w, h = sh['box']
        tb = slide.shapes.add_textbox(Emu(int(x * EMU_PER_PX)), Emu(int(y * EMU_PER_PX)),
                                      Emu(int(w * EMU_PER_PX)), Emu(int(h * EMU_PER_PX)))
        tf = tb.text_frame
        tf.word_wrap = True
        for m in ('margin_left', 'margin_right', 'margin_top', 'margin_bottom'):
            setattr(tf, m, 0)
        try:
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        except Exception:
            pass
        name = pptx_font(sh.get('font', ''))
        bold = bool(sh.get('bold')) or (sh.get('font') in HEAVY)
        hexc = mid_color(sh)
        al = align_map.get((sh.get('align') or '').upper(), PP_ALIGN.LEFT)
        for i, line in enumerate(str(sh['text']).split('\n')):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = al
            _style_run(p.add_run(), sh, name, bold, hexc)
            p.runs[-1].text = line
        n += 1
    return n


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument('page_dirs', nargs='*', help='page folders (each with page.png), in slide order')
    ap.add_argument('--deck', default='', help='auto-discover <deck>/p*/ in page-number order')
    ap.add_argument('-o', '--out', required=True)
    ap.add_argument('--hires', action='store_true', help='re-rasterize each SVG at 2x for crisp slides')
    ap.add_argument('--editable', action='store_true',
                    help='emit NATIVE editable text boxes over a text-free visual background (vs image-only)')
    a = ap.parse_args()

    from pptx import Presentation
    from pptx.util import Inches

    dirs = collect(a)
    if not dirs:
        sys.exit("no page dirs (pass folders or --deck <runs/deck>)")

    prs = Presentation()
    prs.slide_width = Inches(13.333)        # standard 16:9
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    tmp = tempfile.mkdtemp(prefix='pptx_export_')
    n_ok, n_skip = 0, 0
    for d in dirs:
        if a.editable:
            bg = textfree_bg(d, tmp)
            if not bg:
                print(f"  [skip] {os.path.basename(d)} — no page.svg"); n_skip += 1; continue
            slide = prs.slides.add_slide(blank)
            slide.shapes.add_picture(bg, 0, 0, width=prs.slide_width, height=prs.slide_height)
            nt = add_text_boxes(slide, d)
            print(f"  [slide {n_ok+1}] {os.path.basename(d)}  editable ({nt} text boxes)")
        else:
            img = page_image(d, a.hires, tmp)
            if not img:
                print(f"  [skip] {os.path.basename(d)} — no page.png"); n_skip += 1; continue
            slide = prs.slides.add_slide(blank)
            slide.shapes.add_picture(img, 0, 0, width=prs.slide_width, height=prs.slide_height)
            print(f"  [slide {n_ok+1}] {os.path.basename(d)}{'  (2x)' if a.hires else ''}")
        n_ok += 1

    if not n_ok:
        sys.exit("no slides written")
    os.makedirs(os.path.dirname(os.path.abspath(a.out)), exist_ok=True)
    prs.save(a.out)
    mode = 'editable text + visual bg' if a.editable else ('image-only 2x' if a.hires else 'image-only')
    print(f"\nPPTX OK -> {a.out}  ({n_ok} slides{', '+str(n_skip)+' skipped' if n_skip else ''}, 16:9, {mode})")


if __name__ == '__main__':
    main()
