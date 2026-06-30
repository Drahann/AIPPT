#!/usr/bin/env python3
"""parse_chrome — dump + extract a slide's LAYOUT/MASTER persistent chrome.

The full-bleed background, top nav frame, and other persistent decorations often
live in the slide layout / master, NOT in the slide's own shapes (replication
trap #3). This extracts those images + the bg blip so the renderer can reuse them.

Usage: python parse_chrome.py <pptx> <slide_1based> <out_dir>
Writes L*/M* images, _chrome.txt, and chrome.json (bg + shape boxes).
"""
import os, sys, json
from pptx import Presentation

A = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
P = '{http://schemas.openxmlformats.org/presentationml/2006/main}'
R = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}'

src, sno, outdir = sys.argv[1], int(sys.argv[2]), sys.argv[3]
os.makedirs(outdir, exist_ok=True)
prs = Presentation(src)
W, H = int(prs.slide_width), int(prs.slide_height)
SX, SY = 1280.0 / W, 720.0 / H
slide = prs.slides[sno - 1]
layout = slide.slide_layout
master = layout.slide_master
lines, items = [], []


def grp_xfrm(el):
    xf = el.find(P + 'grpSpPr/' + A + 'xfrm')
    o, e = xf.find(A + 'off'), xf.find(A + 'ext')
    co, ce = xf.find(A + 'chOff'), xf.find(A + 'chExt')
    return (int(o.get('x')), int(o.get('y')), int(e.get('cx')), int(e.get('cy')),
            int(co.get('x')), int(co.get('y')), int(ce.get('cx')), int(ce.get('cy')))


def resolve_bg_blip(part, prs):
    """Find a bg blipFill and write the embedded image. Returns filename or None."""
    bg = part._element.find(P + 'cSld/' + P + 'bg')
    if bg is None:
        return None
    blip = bg.find('.//' + A + 'blip')
    if blip is None:
        return None
    rid = blip.get(R + 'embed')
    if not rid:
        return None
    try:
        img = part.part.related_part(rid)
        ext = img.partname.ext if hasattr(img.partname, 'ext') else 'png'
        fn = f"bg_{type(part).__name__}.{ext}"
        open(os.path.join(outdir, fn), 'wb').write(img.blob)
        return fn
    except Exception as e:
        return f"ERR:{str(e)[:30]}"


def walk(shapes, part_label, mapxy, kx, ky, imgn):
    for sh in shapes:
        try:
            l, t, w, h = sh.left, sh.top, sh.width, sh.height
        except Exception:
            l = None
        if sh.shape_type == 6:
            ox, oy, ex, ey, cox, coy, cex, cey = grp_xfrm(sh._element)
            gx0, gy0 = mapxy(ox, oy)
            walk(sh.shapes, part_label,
                 (lambda x, y, gx0=gx0, gy0=gy0, cox=cox, coy=coy, nx=kx*(ex/cex if cex else 1), ny=ky*(ey/cey if cey else 1): (gx0+(x-cox)*nx, gy0+(y-coy)*ny)),
                 kx*(ex/cex if cex else 1), ky*(ey/cey if cey else 1), imgn)
            continue
        if l is None:
            continue
        X, Y = mapxy(l, t)
        box = [round(X*SX, 1), round(Y*SY, 1), round(w*kx*SX, 1), round(h*ky*SY, 1)]
        extra, fn = '', None
        if sh.shape_type == 13:
            try:
                img = sh.image
                imgn[0] += 1
                fn = f"{part_label}{imgn[0]:02d}.{img.ext}"
                open(os.path.join(outdir, fn), 'wb').write(img.blob)
                extra = f" img={fn}({img.size[0]}x{img.size[1]})"
            except Exception as e:
                extra = f" img_err={str(e)[:30]}"
        try:
            if sh.has_text_frame and sh.text_frame.text.strip():
                extra += f"  <<{sh.text_frame.text.strip()[:40]}>>"
        except Exception:
            pass
        lines.append(f"  {part_label} [{box[0]:>6.0f},{box[1]:>4.0f} {box[2]:>5.0f}x{box[3]:>4.0f}] {sh.name[:18]:<18}{extra}")
        if fn:
            items.append({'part': part_label, 'asset': fn, 'box': box, 'name': sh.name})


bg = {}
for part, label in [(slide, 'SLIDE'), (layout, 'LAYOUT'), (master, 'MASTER')]:
    b = resolve_bg_blip(part, prs)
    if b:
        bg[label] = b
        lines.append(f"[BG-{label}] -> {b}")

for shapes, label in [(layout.shapes, 'L'), (master.shapes, 'M')]:
    lines.append(f"--- {label} shapes ---")
    walk(shapes, label, (lambda x, y: (x, y)), 1.0, 1.0, [0])

rep = "\n".join(lines)
open(os.path.join(outdir, '_chrome.txt'), 'w', encoding='utf-8').write(rep)
json.dump({'canvas': [1280, 720], 'bg': bg, 'shapes': items},
          open(os.path.join(outdir, 'chrome.json'), 'w', encoding='utf-8'), ensure_ascii=False, indent=1)
try:
    print(rep)
except Exception:
    print("chrome ->", outdir)
