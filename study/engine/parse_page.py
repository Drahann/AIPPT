#!/usr/bin/env python3
"""parse_page — extract ONE slide's full geometry into a clean page_record (raw layer).

Improves on parse_slide.py: adds srcRect crop, text gradFill (the "name is a gold
gradient, not the red you see" trap), blip alpha/duotone, stable shape ids, and a
structured JSON the deterministic renderer consumes. Roles/slot-ids are added in a
later (VLM/orchestrator) pass — this layer is pure truth.

Usage: python parse_page.py <pptx> <slide_1based> <out_dir>
Writes <out_dir>/page_record.raw.json, _report.txt, and extracts pic*/ images.
"""
import os, sys, json, zipfile
import xml.etree.ElementTree as ET
from pptx import Presentation

A = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
P = '{http://schemas.openxmlformats.org/presentationml/2006/main}'

src, sno, outdir = sys.argv[1], int(sys.argv[2]), sys.argv[3]
os.makedirs(outdir, exist_ok=True)
prs = Presentation(src)

# ---- theme palette + clrMap, for resolving schemeClr (trap #2: colors live in XML) ----
_z = zipfile.ZipFile(src)
_theme = ET.fromstring(_z.read('ppt/theme/theme1.xml').decode('utf-8')).find('.//' + A + 'clrScheme')
THEME = {}
for _c in _theme:
    _t = _c.tag.split('}')[-1]
    _s, _y = _c.find(A + 'srgbClr'), _c.find(A + 'sysClr')
    THEME[_t] = (_s.get('val') if _s is not None else _y.get('lastClr')).upper()
# clrMap from master (bg1/tx1/bg2/tx2 -> scheme slot); default mapping if absent
CLRMAP = {'bg1': 'lt1', 'tx1': 'dk1', 'bg2': 'lt2', 'tx2': 'dk2',
          'dk1': 'dk1', 'lt1': 'lt1', 'dk2': 'dk2', 'lt2': 'lt2',
          'accent1': 'accent1', 'accent2': 'accent2', 'accent3': 'accent3',
          'accent4': 'accent4', 'accent5': 'accent5', 'accent6': 'accent6',
          'hlink': 'hlink', 'folHlink': 'folHlink'}
try:
    _cm = prs.slides[sno - 1].slide_layout.slide_master._element.find(P + 'clrMap')
    if _cm is not None:
        for _k in ('bg1', 'tx1', 'bg2', 'tx2'):
            if _cm.get(_k):
                CLRMAP[_k] = _cm.get(_k)
except Exception:
    pass


def _hex2rgb(h):
    return tuple(int(h[i:i + 2], 16) for i in (0, 2, 4))


def _rgb2hex(r, g, b):
    return '#%02X%02X%02X' % (max(0, min(255, int(r))), max(0, min(255, int(g))), max(0, min(255, int(b))))


def _apply_mods(hexv, clr_el):
    """apply lumMod/lumOff/shade/tint children to a base hex (the gold-name/gray-nav trap)."""
    r, g, b = _hex2rgb(hexv)
    def pct(tag):
        e = clr_el.find(A + tag)
        return int(e.get('val')) / 100000.0 if e is not None and e.get('val') else None
    lm, lo, sh, ti = pct('lumMod'), pct('lumOff'), pct('shade'), pct('tint')
    if lm is not None:
        r, g, b = r * lm, g * lm, b * lm
    if lo is not None:
        r, g, b = r + 255 * lo, g + 255 * lo, b + 255 * lo
    if sh is not None:
        r, g, b = r * sh, g * sh, b * sh
    if ti is not None:
        r, g, b = r + (255 - r) * (1 - ti), g + (255 - g) * (1 - ti), b + (255 - b) * (1 - ti)
    return _rgb2hex(r, g, b)


def _resolve_clr(clr_el):
    """srgbClr or schemeClr element -> #RRGGBB with luminance mods applied."""
    if clr_el is None:
        return None
    tag = clr_el.tag.split('}')[-1]
    if tag == 'srgbClr':
        return _apply_mods(clr_el.get('val').upper(), clr_el)
    if tag == 'schemeClr':
        slot = CLRMAP.get(clr_el.get('val'), clr_el.get('val'))
        base = THEME.get(slot)
        return _apply_mods(base, clr_el) if base else None
    if tag == 'sysClr':
        return '#' + clr_el.get('lastClr', '000000').upper()
    return None
W, H = int(prs.slide_width), int(prs.slide_height)
SX, SY = 1280.0 / W, 720.0 / H
slide = prs.slides[sno - 1]
rows, imgn = [], [0]


def grp_xfrm(el):
    xf = el.find(P + 'grpSpPr/' + A + 'xfrm')
    o, e = xf.find(A + 'off'), xf.find(A + 'ext')
    co, ce = xf.find(A + 'chOff'), xf.find(A + 'chExt')
    return (int(o.get('x')), int(o.get('y')), int(e.get('cx')), int(e.get('cy')),
            int(co.get('x')), int(co.get('y')), int(ce.get('cx')), int(ce.get('cy')))


def color_of(fill_el):
    """solidFill element -> #RRGGBB, resolving srgbClr/schemeClr + luminance mods."""
    if fill_el is None:
        return None
    for tag in ('srgbClr', 'schemeClr', 'sysClr'):
        c = fill_el.find(A + tag)
        if c is not None:
            return _resolve_clr(c)
    return None


def grad_of(rpr):
    """Text run gradFill -> list of stop colors (the gold-name trap)."""
    g = rpr.find(A + 'gradFill') if rpr is not None else None
    if g is None:
        return None
    stops = []
    for gs in g.findall('.//' + A + 'gs'):
        for tag in ('srgbClr', 'schemeClr', 'sysClr'):
            c = gs.find(A + tag)
            if c is not None:
                stops.append(_resolve_clr(c)); break
    return stops or None


def lststyle_color(shape_el):
    """Default run color inherited from the shape's txBody lstStyle (inactive nav etc.)."""
    ls = shape_el.find('.//' + A + 'lstStyle')
    if ls is None:
        return None
    drpr = ls.find('.//' + A + 'defRPr')
    if drpr is None:
        return None
    return color_of(drpr.find(A + 'solidFill'))


def srcrect_of(pic_el):
    sr = pic_el.find('.//' + A + 'blipFill/' + A + 'srcRect')
    if sr is None:
        return None
    d = {}
    for k in ('l', 't', 'r', 'b'):
        v = sr.get(k)
        if v:
            d[k] = round(int(v) / 100000.0, 4)   # ST_Percentage 100000 == 100%
    return d or None


def blip_fx(pic_el):
    """alpha + duotone baked-render hints living on the blip."""
    blip = pic_el.find('.//' + A + 'blip')
    if blip is None:
        return None, None
    am = blip.find(A + 'alphaModFix')
    alpha = round(int(am.get('amt')) / 100000.0, 3) if am is not None and am.get('amt') else None
    duo = blip.find(A + 'duotone')
    duocols = [('#' + c.get('val')) for c in duo.findall(A + 'srgbClr')] if duo is not None else None
    # prstClr (e.g. white/black) in duotone
    if duo is not None and not duocols:
        duocols = [c.get('val') for c in duo.findall(A + 'prstClr')]
    return alpha, (duocols or None)


# ============ EXHAUSTIVE detail extractors (geom / fill / line / effects / text) ============
EMU_PX = 9525.0   # 1 px @ 96dpi


def _px(emu):
    try:
        return round(int(emu) / EMU_PX, 1)
    except Exception:
        return None


def _alpha(clr_el):
    a = clr_el.find(A + 'alpha') if clr_el is not None else None
    return round(int(a.get('val')) / 100000.0, 3) if a is not None and a.get('val') else None


def geom_of(spPr):
    g = spPr.find(A + 'prstGeom')
    if g is not None:
        d = {'type': g.get('prst')}
        adj = [av.get('fmla') for av in g.findall(A + 'avLst/' + A + 'gd')]
        if adj:
            d['adj'] = adj
        return d
    cg = spPr.find(A + 'custGeom')
    if cg is not None:
        d = {'type': 'custom'}
        paths = []
        try:
            for path in cg.findall(A + 'pathLst/' + A + 'path'):
                pw, ph = float(path.get('w') or 0), float(path.get('h') or 0)
                if not (pw and ph):
                    continue
                cmds = []
                for el in path:
                    tag = el.tag.split('}')[-1]
                    pts = [(float(p.get('x')) / pw, float(p.get('y')) / ph) for p in el.findall(A + 'pt')]
                    if tag == 'moveTo' and pts:
                        cmds.append('M%.4f,%.4f' % pts[0])
                    elif tag == 'lnTo' and pts:
                        cmds.append('L%.4f,%.4f' % pts[0])
                    elif tag == 'cubicBezTo' and len(pts) == 3:
                        cmds.append('C%.4f,%.4f %.4f,%.4f %.4f,%.4f' % (pts[0][0], pts[0][1], pts[1][0], pts[1][1], pts[2][0], pts[2][1]))
                    elif tag == 'quadBezTo' and len(pts) == 2:
                        cmds.append('Q%.4f,%.4f %.4f,%.4f' % (pts[0][0], pts[0][1], pts[1][0], pts[1][1]))
                    elif tag == 'close':
                        cmds.append('Z')
                if cmds:
                    paths.append(''.join(cmds))      # normalized 0..1 -> scale to box to re-render (recolorable!)
        except Exception:
            pass
        if paths:
            d['path'] = paths
        return d
    return None


def grad_fill_of(gf):
    stops = []
    for gs in gf.findall(A + 'gsLst/' + A + 'gs'):
        pos = round(int(gs.get('pos', '0')) / 100000.0, 3)
        for tag in ('srgbClr', 'schemeClr', 'sysClr'):
            c = gs.find(A + tag)
            if c is not None:
                stops.append({'pos': pos, 'color': _resolve_clr(c), 'alpha': _alpha(c)})
                break
    lin = gf.find(A + 'lin')
    ang = round(int(lin.get('ang')) / 60000.0, 1) if lin is not None and lin.get('ang') else None
    return {'type': 'grad', 'stops': stops, 'angle': ang}


def fill_of(spPr):
    if spPr is None:
        return None
    if spPr.find(A + 'noFill') is not None:
        return {'type': 'none'}
    sf = spPr.find(A + 'solidFill')
    if sf is not None:
        for tag in ('srgbClr', 'schemeClr', 'sysClr'):
            c = sf.find(A + tag)
            if c is not None:
                return {'type': 'solid', 'color': _resolve_clr(c), 'alpha': _alpha(c)}
    gf = spPr.find(A + 'gradFill')
    if gf is not None:
        return grad_fill_of(gf)
    if spPr.find(A + 'blipFill') is not None:
        return {'type': 'pic'}
    if spPr.find(A + 'pattFill') is not None:
        return {'type': 'pattern'}
    return None


def line_of(spPr):
    if spPr is None:
        return None
    ln = spPr.find(A + 'ln')
    if ln is None:
        return None
    if ln.find(A + 'noFill') is not None:
        return {'none': True}
    d = {}
    if ln.get('w'):
        d['w'] = _px(ln.get('w'))
    sf = ln.find(A + 'solidFill')
    if sf is not None:
        for tag in ('srgbClr', 'schemeClr', 'sysClr'):
            c = sf.find(A + tag)
            if c is not None:
                d['color'] = _resolve_clr(c)
                a = _alpha(c)
                if a is not None:
                    d['alpha'] = a
                break
    gf = ln.find(A + 'gradFill')
    if gf is not None:
        d['grad'] = grad_fill_of(gf)
    dash = ln.find(A + 'prstDash')
    if dash is not None:
        d['dash'] = dash.get('val')
    return d or None


def effects_of(spPr):
    if spPr is None:
        return None
    lst = spPr.find(A + 'effectLst')
    if lst is None:
        return None
    d = {}
    sh = lst.find(A + 'outerShdw')
    if sh is not None:
        s = {}
        for k, attr in (('blur', 'blurRad'), ('dist', 'dist')):
            if sh.get(attr):
                s[k] = _px(sh.get(attr))
        if sh.get('dir'):
            s['dir'] = round(int(sh.get('dir')) / 60000.0, 1)
        c = sh.find(A + 'srgbClr')
        c = c if c is not None else sh.find(A + 'schemeClr')
        if c is not None:
            s['color'] = _resolve_clr(c)
            a = _alpha(c)
            if a is not None:
                s['alpha'] = a
        d['shadow'] = s
    gl = lst.find(A + 'glow')
    if gl is not None:
        g = {}
        if gl.get('rad'):
            g['rad'] = _px(gl.get('rad'))
        c = gl.find(A + 'srgbClr')
        c = c if c is not None else gl.find(A + 'schemeClr')
        if c is not None:
            g['color'] = _resolve_clr(c)
        d['glow'] = g
    se = lst.find(A + 'softEdge')
    if se is not None and se.get('rad'):
        d['softEdge'] = _px(se.get('rad'))
    if lst.find(A + 'reflection') is not None:
        d['reflection'] = True
    return d or None


_ALGN = {'l': 'LEFT', 'ctr': 'CENTER', 'r': 'RIGHT', 'just': 'JUSTIFY', 'dist': 'DISTRIBUTE'}


def textframe_of(tb):
    bp = tb.find(A + 'bodyPr')
    tf = {}
    if bp is not None:
        if bp.get('anchor'):
            tf['anchor'] = bp.get('anchor')        # t/ctr/b
        ins = {}
        for k, attr in (('l', 'lIns'), ('t', 'tIns'), ('r', 'rIns'), ('b', 'bIns')):
            if bp.get(attr) is not None:
                ins[k] = _px(bp.get(attr))
        if ins:
            tf['ins'] = ins
        if bp.get('wrap'):
            tf['wrap'] = bp.get('wrap')
        if bp.find(A + 'normAutofit') is not None:
            tf['autofit'] = 'norm'
        elif bp.find(A + 'spAutoFit') is not None:
            tf['autofit'] = 'shape'
        elif bp.find(A + 'noAutofit') is not None:
            tf['autofit'] = 'none'
        if bp.get('vert'):
            tf['vert'] = bp.get('vert')             # vertical text (eaVert etc.)
    return tf or None


def run_props(r):
    d = {}
    te = r.find(A + 't')
    d['t'] = te.text if (te is not None and te.text) else ''
    rpr = r.find(A + 'rPr')
    if rpr is not None:
        if rpr.get('sz'):
            pt = int(rpr.get('sz')) / 100.0
            d['font_pt'] = round(pt, 1)
            d['font_px'] = round(pt * 4 / 3, 1)
        for tag in ('latin', 'ea', 'cs'):
            e = rpr.find(A + tag)
            if e is not None and e.get('typeface'):
                d['font'] = e.get('typeface')
                break
        if rpr.get('b') == '1':
            d['bold'] = True
        if rpr.get('i') == '1':
            d['italic'] = True
        if rpr.get('u') and rpr.get('u') != 'none':
            d['underline'] = rpr.get('u')
        if rpr.get('strike') and rpr.get('strike') != 'noStrike':
            d['strike'] = rpr.get('strike')
        if rpr.get('spc'):
            d['spc'] = round(int(rpr.get('spc')) / 100.0, 1)
        if rpr.get('baseline'):
            d['baseline'] = round(int(rpr.get('baseline')) / 1000.0, 2)
        g = grad_of(rpr)
        if g:
            d['grad'] = g
        else:
            sol = color_of(rpr.find(A + 'solidFill'))
            if sol:
                d['color'] = sol
        ln = rpr.find(A + 'ln')   # text outline
        if ln is not None:
            lc = color_of(ln.find(A + 'solidFill'))
            if lc:
                d['outline'] = {'color': lc, 'w': _px(ln.get('w')) if ln.get('w') else None}
    return d


def para_props(p):
    d = {}
    ppr = p.find(A + 'pPr')
    if ppr is not None:
        if ppr.get('algn'):
            d['align'] = _ALGN.get(ppr.get('algn'), ppr.get('algn'))
        if ppr.get('lvl'):
            d['lvl'] = int(ppr.get('lvl'))
        for attr, key in (('marL', 'marL'), ('indent', 'indent')):
            if ppr.get(attr):
                d[key] = _px(ppr.get(attr))
        ls = ppr.find(A + 'lnSpc')
        if ls is not None:
            pct, pts = ls.find(A + 'spcPct'), ls.find(A + 'spcPts')
            if pct is not None:
                d['lnSpc'] = round(int(pct.get('val')) / 100000.0, 2)
            elif pts is not None:
                d['lnSpcPt'] = round(int(pts.get('val')) / 100.0, 1)
        for tag, key in (('spcBef', 'spcBef'), ('spcAft', 'spcAft')):
            e = ppr.find(A + tag)
            if e is not None and e.find(A + 'spcPts') is not None:
                d[key] = round(int(e.find(A + 'spcPts').get('val')) / 100.0, 1)
        if ppr.find(A + 'buNone') is not None:
            d['bullet'] = 'none'
        elif ppr.find(A + 'buChar') is not None:
            d['bullet'] = {'char': ppr.find(A + 'buChar').get('char')}
        elif ppr.find(A + 'buAutoNum') is not None:
            d['bullet'] = {'auto': ppr.find(A + 'buAutoNum').get('type')}
    runs = [run_props(r) for r in p.findall(A + 'r')]
    runs = [r for r in runs if r.get('t')]
    if runs:
        d['runs'] = runs
    return d


def walk(shapes, mapxy, kx, ky, zbase=0):
    for sh in shapes:
        try:
            l, t, w, h = sh.left, sh.top, sh.width, sh.height
        except Exception:
            l = t = w = h = None
        if sh.shape_type == 6:  # GROUP
            ox, oy, ex, ey, cox, coy, cex, cey = grp_xfrm(sh._element)
            gx0, gy0 = mapxy(ox, oy)
            nkx, nky = kx * (ex / cex if cex else 1), ky * (ey / cey if cey else 1)
            walk(sh.shapes,
                 (lambda x, y, gx0=gx0, gy0=gy0, cox=cox, coy=coy, nkx=nkx, nky=nky:
                  (gx0 + (x - cox) * nkx, gy0 + (y - coy) * nky)), nkx, nky)
            continue
        if l is None:
            continue
        X, Y = mapxy(l, t)
        box = [round(X * SX, 1), round(Y * SY, 1), round(w * kx * SX, 1), round(h * ky * SY, 1)]
        sid = f"sh{len(rows):02d}"
        rec = {'id': sid, 'name': sh.name, 'kind': 'shape', 'box': box, 'z': len(rows)}
        xf = sh._element.find('.//' + A + 'xfrm')
        if xf is not None:
            rot = int(xf.get('rot', '0'))
            if rot:
                rec['rot'] = round(rot / 60000.0, 1)
            if xf.get('flipH') in ('1', 'true'):
                rec['flipH'] = True
            if xf.get('flipV') in ('1', 'true'):
                rec['flipV'] = True
        # geometry / fill / outline / effects (full detail)
        spPr = sh._element.find('.//' + P + 'spPr')
        if spPr is not None:
            g = geom_of(spPr)
            if g:
                rec['geom'] = g
            fl = fill_of(spPr)
            if fl:
                rec['fillp'] = fl
                if fl.get('type') == 'solid' and fl.get('color'):
                    rec['fill'] = fl['color']        # back-compat flat color
            ln = line_of(spPr)
            if ln:
                rec['line'] = ln
            ef = effects_of(spPr)
            if ef:
                rec['effects'] = ef
        # ---- text: ALL paragraphs + ALL runs, full props ----
        try:
            if sh.has_text_frame and sh.text_frame.text.strip():
                rec['kind'] = 'text'
                tb = sh.text_frame._txBody
                rec['text'] = sh.text_frame.text.strip()
                tfp = textframe_of(tb)
                if tfp:
                    rec['tf'] = tfp
                paras = [para_props(p) for p in tb.findall(A + 'p')]
                paras = [p for p in paras if p.get('runs')]
                if paras:
                    rec['paras'] = paras
                    r0 = paras[0]['runs'][0]            # back-compat flat fields (first run)
                    for k in ('font', 'font_px', 'font_pt', 'bold', 'italic', 'color', 'grad'):
                        if k in r0:
                            rec[k] = r0[k]
                    if paras[0].get('align'):
                        rec['align'] = paras[0]['align']
                if 'color' not in rec and 'grad' not in rec:
                    inh = lststyle_color(sh._element)
                    if inh:
                        rec['color'] = inh
        except Exception:
            pass
        # ---- picture ----
        if sh.shape_type == 13:
            try:
                img = sh.image
                imgn[0] += 1
                fn = f"pic{imgn[0]:02d}.{img.ext}"
                open(os.path.join(outdir, fn), 'wb').write(img.blob)
                rec['kind'] = 'pic'
                rec['image'] = fn
                rec['img_native'] = f"{img.size[0]}x{img.size[1]}"
                sr = srcrect_of(sh._element)
                if sr:
                    rec['srcRect'] = sr
                alpha, duo = blip_fx(sh._element)
                if alpha is not None:
                    rec['alpha'] = alpha
                if duo:
                    rec['duotone'] = duo
            except Exception as e:
                rec['image_err'] = str(e)[:40]
        rows.append(rec)


walk(slide.shapes, (lambda x, y: (x, y)), 1.0, 1.0)

# human report
lines = [f"slide {sno}/{len(prs.slides)} | EMU {W}x{H} -> 1280x720 | shapes:{len(rows)} imgs:{imgn[0]}", "=" * 110]
for r in rows:
    bx = r['box']
    tag = {'pic': 'PIC', 'text': 'TXT', 'shape': 'shp'}[r['kind']]
    fl = ''.join(f" {k}{r[k]}" if k == 'rot' else f" {k}" for k in ('rot', 'flipH', 'flipV') if r.get(k))
    extra = ''
    if r['kind'] == 'pic':
        extra = f" {r['image']}({r['img_native']})"
        if 'srcRect' in r:
            extra += f" crop={r['srcRect']}"
        if 'duotone' in r:
            extra += f" duo={r['duotone']}"
        if 'alpha' in r:
            extra += f" a={r['alpha']}"
    if r.get('fill'):
        extra += f" fill={r['fill']}"
    if 'text' in r:
        f = []
        if 'font' in r: f.append(r['font'])
        if 'font_px' in r: f.append(f"{r['font_px']}px")
        if r.get('bold'): f.append('B')
        if 'grad' in r: f.append('grad' + str(r['grad']))
        elif 'color' in r: f.append(r['color'])
        if 'align' in r: f.append(r['align'])
        t = r['text'].replace('\n', ' / ')[:50]
        extra += f"  {{{' '.join(f)}}}  <<{t}>>"
    lines.append(f"{r['id']} {tag}[{bx[0]:>6.0f},{bx[1]:>4.0f} {bx[2]:>5.0f}x{bx[3]:>4.0f}]{fl}{extra}")
report = "\n".join(lines)
open(os.path.join(outdir, '_report.txt'), 'w', encoding='utf-8').write(report)
json.dump({'canvas': [1280, 720], 'source': {'pptx': os.path.basename(src), 'slide': sno}, 'shapes': rows},
          open(os.path.join(outdir, 'page_record.raw.json'), 'w', encoding='utf-8'), ensure_ascii=False, indent=1)
try:
    print(report)
except Exception:
    print(f"parsed {len(rows)} shapes / {imgn[0]} imgs -> {outdir}")
